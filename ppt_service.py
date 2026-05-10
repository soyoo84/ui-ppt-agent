import os
import json
import logging
import datetime
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from schemas import ScreenAnalysisResult
from config import (
    MASTER_PPT_PATH, TARGET_LAYOUT_NAME, PPT_SLIDE_WIDTH, PPT_SLIDE_HEIGHT, PPT_UI_SCALE,
    PPT_TOP_OFFSET_RATIO, PPT_ALIGN_THRESHOLD, PPT_CONTAINER_PADDING, HDS_PRIMARY_COLOR
)
from prompts import get_component_registry

logger = logging.getLogger(__name__)

# --- [유틸리티 헬퍼 함수 모음] ---
def _safe_rgb(val, default_rgb):
    """[안전 방어] JSON 휴먼 에러(타입 오타, 값 누락)를 방어하는 RGB 헬퍼 함수"""
    if isinstance(val, list) and len(val) >= 3:
        try:
            return RGBColor(max(0, min(255, int(val[0]))), max(0, min(255, int(val[1]))), max(0, min(255, int(val[2]))))
        except Exception: pass
    return RGBColor(*default_rgb)

def _hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip('#')
    if len(hex_color) == 3: hex_color = ''.join([c*2 for c in hex_color])
    if len(hex_color) == 6: return [int(hex_color[i:i+2], 16) for i in (0, 2, 4)]
    return None

def _parse_color(color_str):
    color_str = color_str.strip().lower()
    if color_str.startswith('#'): return _hex_to_rgb(color_str)
    elif color_str.startswith('rgb'):
        match = re.search(r'rgba?\((\d+),\s*(\d+),\s*(\d+)', color_str)
        if match: return [int(match.group(1)), int(match.group(2)), int(match.group(3))]
    return None

def get_layout_names(template_path):
    """PPT 템플릿 파일에서 슬라이드 레이아웃 이름 목록을 추출합니다."""
    if not template_path or not os.path.exists(template_path):
        return []
    try:
        prs = Presentation(template_path)
        return [layout.name for layout in prs.slide_layouts]
    except Exception:
        return []

def create_editable_ppt(analysis_result: ScreenAnalysisResult, output_file, template_path=None, layout_name=None):
    """
    분석된 JSON 데이터를 바탕으로 수정 가능한 PPT 파일을 생성합니다.
    """
    # 사내 템플릿 로드 (파일이 존재하지 않으면 기본 빈 프레젠테이션 생성)
    if template_path and os.path.exists(template_path):
        prs = Presentation(template_path)
    # 사용자가 아직 master 폴더로 이관하지 않고 기존 환경변수를 사용하는 경우에 대한 안전한 롤백(호환성)
    elif MASTER_PPT_PATH and os.path.exists(MASTER_PPT_PATH):
        prs = Presentation(MASTER_PPT_PATH)
    else:
        prs = Presentation()
        # 16:9 와이드스크린 비율 설정
        prs.slide_width = Inches(PPT_SLIDE_WIDTH)
        prs.slide_height = Inches(PPT_SLIDE_HEIGHT)
    
    # --- [레이아웃 이름으로 템플릿 찾기] ---
    # 사용하는 템플릿의 슬라이드 마스터에 있는 실제 레이아웃 이름을 입력하세요.
    target_layout_name = layout_name or TARGET_LAYOUT_NAME 
    slide_layout = None
    
    if target_layout_name:
        for layout in prs.slide_layouts:
            if layout.name == target_layout_name:
                slide_layout = layout
                break
            
    # 설정한 이름의 레이아웃을 찾지 못한 경우 (Fallback)
    if slide_layout is None:
        slide_layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0]
    
    slide = prs.slides.add_slide(slide_layout)
    
    # 제목 설정
    if slide.shapes.title:
        slide.shapes.title.text = analysis_result.screen_name
    
    # 슬라이드 전체 너비/높이 (원본)
    actual_slide_width = prs.slide_width
    actual_slide_height = prs.slide_height
    
    # 템플릿 상단의 'UI 정의서 표/내용' 영역을 덮어쓰지 않도록 UI 렌더링 영역을 하단으로 내립니다.
    ui_scale = PPT_UI_SCALE
    # 상단 여백 확보를 위해 UI 스케일이 너무 크면 자동으로 약간 축소 (기본 65%)
    if ui_scale > 0.65:
        ui_scale = 0.65
        
    slide_width = actual_slide_width * ui_scale
    slide_height = actual_slide_height * ui_scale
    
    # UI 영역 위치 오프셋 (좌측 여백 5%, 설정된 상단 여백 비율을 적용하여 하단 좌측으로 배치)
    offset_left = int(actual_slide_width * 0.05)
    offset_top = int(actual_slide_height * PPT_TOP_OFFSET_RATIO)
    
    # --- [정렬 보정 알고리즘 추가 시작] ---
    # 설정된 오차 범위 이내인 요소들을 같은 그룹(행)으로 묶습니다.
    y_threshold = PPT_ALIGN_THRESHOLD 
    sorted_comps = sorted(analysis_result.components, key=lambda c: c.y_percent)
    
    if sorted_comps:
        groups, current_group = [], [sorted_comps[0]]
        for comp in sorted_comps[1:]:
            if abs(comp.y_percent - current_group[0].y_percent) <= y_threshold:
                current_group.append(comp)
            else:
                groups.append(current_group)
                current_group = [comp]
        groups.append(current_group)
        
        # 같은 행에 있는 요소들의 Y 좌표를 그룹의 평균값으로 통일하여 수평 정렬을 맞춥니다. (높이는 원본 유지)
        for group in groups:
            avg_y = sum(c.y_percent for c in group) / len(group)
            for c in group:
                c.y_percent = avg_y
    # --- [정렬 보정 알고리즘 추가 끝] ---
    
    # --- [수직 정렬(Column Alignment) 보정 알고리즘 추가 시작] ---
    # 설정된 오차 범위 이내인 요소들을 같은 그룹(열)으로 묶습니다.
    x_threshold = PPT_ALIGN_THRESHOLD 
    sorted_comps_x = sorted(analysis_result.components, key=lambda c: c.x_percent)
    
    if sorted_comps_x:
        groups_x, current_group_x = [], [sorted_comps_x[0]]
        for comp in sorted_comps_x[1:]:
            if abs(comp.x_percent - current_group_x[0].x_percent) <= x_threshold:
                current_group_x.append(comp)
            else:
                groups_x.append(current_group_x)
                current_group_x = [comp]
        groups_x.append(current_group_x)
        
        # 같은 열에 있는 요소들의 X 좌표를 그룹의 평균값으로 통일하여 왼쪽 맞춤을 합니다.
        for group in groups_x:
            avg_x = sum(c.x_percent for c in group) / len(group)
            for c in group:
                c.x_percent = avg_x
    # --- [수직 정렬(Column Alignment) 보정 알고리즘 추가 끝] ---

    # --- [시각적 컨테이너(배경 화면 박스) 추가 시작] ---
    if analysis_result.components:
        min_x = min(c.x_percent for c in analysis_result.components)
        min_y = min(c.y_percent for c in analysis_result.components)
        max_x = max(c.x_percent + c.width_percent for c in analysis_result.components)
        max_y = max(c.y_percent + c.height_percent for c in analysis_result.components)
        
        # 상하좌우 여백(Padding) 3% 추가
        pad = PPT_CONTAINER_PADDING
        bg_left = int(slide_width * max(0, min_x - pad)) + offset_left
        bg_top = int(slide_height * max(0, min_y - pad)) + offset_top
        bg_width = int(slide_width * min(1.0, (max_x - min_x + pad * 2)))
        bg_height = int(slide_height * min(1.0, (max_y - min_y + pad * 2)))
        
        # 둥근 사각형 배경 패널 그리기 (도형을 먼저 그리면 자동으로 맨 밑(Z-index 최하단)에 깔립니다)
        bg_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bg_left, bg_top, bg_width, bg_height)
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = RGBColor(248, 249, 250) # 옅은 배경색 (Light Gray)
        bg_shape.line.color.rgb = RGBColor(220, 220, 220)      # 옅은 테두리선
    # --- [시각적 컨테이너(배경 화면 박스) 추가 끝] ---

    # --- [Z-Index 보정 알고리즘 추가] ---
    # 모달, 툴팁, 드롭다운 등 화면에 떠 있는(Floating) 요소들이 다른 도형에 가려지지 않도록 가장 마지막(최상단)에 그립니다.
    floating_types = {"Modal", "HdsModal", "Dialog", "ant-modal", "Tooltip", "HdsTooltip", "Popover", "ant-tooltip", "ant-tooltip-inner", "Dropdown", "Select", "HdsSelect", "HdsDropdown", "ant-select", "ant-dropdown"}
    regular_comps = [c for c in analysis_result.components if c.component_type not in floating_types]
    floating_comps = [c for c in analysis_result.components if c.component_type in floating_types]
    
    # 반환된 컴포넌트들을 순회하며 그리기 (일반 도형 먼저 -> 플로팅 도형 나중에)
    for comp in regular_comps + floating_comps:
        # 비율(0.0~1.0)을 안전하게 클램핑(Clamping)하여 PPT 영역을 벗어나지 않도록 방어
        safe_x = max(0.0, min(1.0, comp.x_percent))
        safe_y = max(0.0, min(1.0, comp.y_percent))
        safe_w = max(0.0, min(1.0 - safe_x, comp.width_percent)) # 우측 화면 밖으로 나가는 것 방지
        safe_h = max(0.0, min(1.0 - safe_y, comp.height_percent)) # 하단 화면 밖으로 나가는 것 방지
        
        # PPT의 실제 길이(EMU 단위)로 변환 (오프셋 반영)
        left = int(slide_width * safe_x) + offset_left
        top = int(slide_height * safe_y) + offset_top
        width = int(slide_width * safe_w)
        height = int(slide_height * safe_h)
        
        # 최소 크기 지정 (오류값 보정)
        width = max(int(Inches(0.5)), width)
        height = max(int(Inches(0.3)), height)

        if comp.component_type in ("PrimaryButton", "Button", "HdsButton", "ant-btn", "ant-btn-primary", "ant-btn-default"):
            # 둥근 모서리 사각형 (MSO_SHAPE.ROUNDED_RECTANGLE = 5)
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*HDS_PRIMARY_COLOR) # 환경설정 메인 컬러 적용
            shape.line.color.rgb = RGBColor(*HDS_PRIMARY_COLOR)
            
            tf = shape.text_frame
            tf.text = comp.text if comp.text else "버튼"
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.size = Pt(14)
            
        elif comp.component_type in ("Tooltip", "HdsTooltip", "Popover", "ant-tooltip", "ant-tooltip-inner"):
            # Ant Design 스타일의 Tooltip (어두운 배경의 말풍선 형태)
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(60, 60, 60) # 진한 회색/검정색 배경
            shape.line.color.rgb = RGBColor(60, 60, 60)
            
            tf = shape.text_frame
            tf.text = comp.text if comp.text else "도움말 툴팁"
            tf.word_wrap = True
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            tf.paragraphs[0].font.size = Pt(10)
            
        elif comp.component_type in ("ProgressBar", "Progress", "HdsProgress", "HdsProgressBar", "ant-progress", "ant-progress-bg"):
            # Ant Design 스타일의 Progress Bar (배경 트랙 + 메인 컬러 채움)
            bg_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = RGBColor(240, 240, 240) # 옅은 배경 트랙
            bg_shape.line.color.rgb = RGBColor(230, 230, 230)
            
            # 시각적 기획을 위해 50% 정도 채워진 상태를 기본으로 덧그림
            fg_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, int(width * 0.5), height)
            fg_shape.fill.solid()
            fg_shape.fill.fore_color.rgb = RGBColor(*HDS_PRIMARY_COLOR)
            fg_shape.line.color.rgb = RGBColor(*HDS_PRIMARY_COLOR)
            
        elif comp.component_type in ("TextInput", "Input", "HdsInput", "TextArea", "ant-input"):
            # 일반 사각형 테두리 (MSO_SHAPE.RECTANGLE = 1)
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
            shape.line.color.rgb = RGBColor(180, 180, 180) # 회색 테두리
            
            tf = shape.text_frame
            safe_text = comp.text if comp.text else ""
            tf.text = f" {safe_text}" # 좌측 여백을 위해 공백 추가
            tf.word_wrap = True
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT
            tf.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
            
        elif comp.component_type in ("Dropdown", "Select", "HdsSelect", "HdsDropdown", "ant-select", "ant-dropdown"):
            # Ant Design 스타일의 Select/Dropdown (우측에 ▼ 화살표 추가)
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
            shape.line.color.rgb = RGBColor(200, 200, 200)
            
            tf = shape.text_frame
            safe_text = comp.text if comp.text else "선택"
            tf.text = f" {safe_text}   ▼"
            tf.word_wrap = True
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT
            tf.paragraphs[0].font.color.rgb = RGBColor(80, 80, 80)
            tf.paragraphs[0].font.size = Pt(12)
            
        elif comp.component_type in ("AgGrid", "Table", "DataGrid", "DataTable", "Grid", "HdsDataGrid", "ag-root-wrapper", "ant-table"):
            # 데이터 그리드를 확실한 표(Table) 형태로 렌더링
            # 3행 4열의 기본 표를 생성하여 시각적으로 풍부하게 표현
            table_shape = slide.shapes.add_table(3, 4, left, top, width, height)
            table = table_shape.table
            
            # 헤더(첫 행) 및 데이터 행 스타일링
            for r_idx in range(3):
                for c_idx in range(4):
                    cell = table.cell(r_idx, c_idx)
                    if r_idx == 0:
                        cell.text = f"Column {c_idx + 1}"
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(240, 245, 250) # 세련된 연한 파란색 헤더
                        cell.text_frame.paragraphs[0].font.bold = True
                    else:
                        cell.text = comp.text if c_idx == 0 and r_idx == 1 and comp.text else "Data"
                    
                    cell.text_frame.paragraphs[0].font.size = Pt(10)
                    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(80, 80, 80)
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                    
        elif comp.component_type in ("DatePicker", "HdsDatePicker", "Calendar", "ant-picker"):
            # Ant Design 스타일의 DatePicker (입력창 + 우측 달력 아이콘)
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
            shape.line.color.rgb = RGBColor(180, 180, 180)
            
            tf = shape.text_frame
            # 날짜 텍스트가 없으면 기본값 YYYY-MM-DD 삽입
            safe_text = comp.text if comp.text else "YYYY-MM-DD"
            tf.text = f" {safe_text}   📅"
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT
            tf.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
            tf.paragraphs[0].font.size = Pt(12)
            
        elif comp.component_type in ("AgGridToolbar", "TableToolbar", "GridToolbar", "Toolbar", "HdsDataGridToolbar", "ag-header"):
            # AG Grid 상단 툴바 (우측 정렬된 액션 버튼이나 검색창 표현)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            # LLM이 추출한 텍스트가 있으면 넣고, 없으면 기본 툴바 텍스트를 넣습니다.
            tf.text = comp.text if comp.text else "🔍 검색   ⬇ 엑셀 다운로드"
            tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
            tf.paragraphs[0].font.color.rgb = RGBColor(80, 80, 80)
            tf.paragraphs[0].font.size = Pt(11)
            tf.paragraphs[0].font.bold = True
            
        elif comp.component_type in ("AgGridPagination", "TablePagination", "GridPagination", "Pagination", "HdsPagination", "ant-pagination"):
            # AG Grid 하단 페이징 영역 (가운데 정렬된 페이지 번호)
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
            shape.line.color.rgb = RGBColor(220, 220, 220)
            
            tf = shape.text_frame
            tf.text = "〈   1   2   3   4   5   〉"
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            tf.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
            tf.paragraphs[0].font.size = Pt(11)
            
        elif comp.component_type in ("Checkbox", "HdsCheckbox", "ant-checkbox", "ant-checkbox-wrapper", "ant-checkbox-checked"):
            # Checkbox: 텍스트 박스를 그리고 좌측에 체크박스 특수문자(☑ 또는 ☐) 삽입
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            # Ant Design의 체크 안 된 상태의 기본 박스 모양 표현
            safe_text = comp.text if comp.text else "체크박스"
            tf.text = f"☐ {safe_text}" 
            tf.word_wrap = True
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT
            tf.paragraphs[0].font.color.rgb = RGBColor(50, 50, 50)
            tf.paragraphs[0].font.size = Pt(12)
            
        elif comp.component_type in ("ToggleSwitch", "Switch", "HdsSwitch", "ant-switch", "ant-switch-checked"):
            # Ant Design 스타일의 Toggle Switch (활성화 상태 표현)
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height) # 둥근 사각형
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(24, 144, 255) # Ant Design 기본 파란색 (활성화)
            shape.line.color.rgb = RGBColor(24, 144, 255)
            
            tf = shape.text_frame
            tf.text = " O" # 스위치 손잡이(노브)를 텍스트로 단순 표현
            tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
            tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            
        elif comp.component_type in ("RadioButton", "Radio", "HdsRadio", "ant-radio", "ant-radio-wrapper", "ant-radio-checked"):
            # Ant Design 스타일의 Radio Button (동그란 라디오 버튼 + 텍스트)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            safe_text = comp.text if comp.text else "라디오"
            tf.text = f"◉ {safe_text}" # 선택된 라디오 버튼 기호
            tf.word_wrap = True
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT
            tf.paragraphs[0].font.color.rgb = RGBColor(50, 50, 50)
            tf.paragraphs[0].font.size = Pt(12)
            
        elif comp.component_type in ("Modal", "HdsModal", "Dialog", "ant-modal"):
            # Ant Design 스타일의 모달(Modal) 창 (팝업 컨테이너)
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
            shape.line.color.rgb = RGBColor(120, 120, 120) # 팝업 느낌을 주기 위한 진한 테두리
            
            tf = shape.text_frame
            tf.text = f"  {comp.text}" if comp.text else "  모달 타이틀"
            tf.word_wrap = True
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT
            tf.paragraphs[0].font.color.rgb = RGBColor(30, 30, 30)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.size = Pt(14)

        elif comp.component_type in ("TextLabel", "Label", "Text", "Typography"):
            # 배경이 없는 투명한 텍스트 박스
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = comp.text if comp.text else "텍스트"
            tf.word_wrap = True
            tf.paragraphs[0].font.color.rgb = RGBColor(30, 30, 30)
            tf.paragraphs[0].font.size = Pt(14)
            
        elif comp.component_type in ("ImagePlaceholder", "Image", "Picture", "HdsImage"):
            # 대각선이 교차하는 이미지 자리표시자 스타일 (MSO_SHAPE.RECTANGLE = 1)
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
            shape.line.color.rgb = RGBColor(200, 200, 200)
            
            tf = shape.text_frame
            tf.text = "[ 이미지 영역 ]"
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            tf.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
            
        else:
            # [자동화 폴백] 외부 JSON 통합 레지스트리를 스캔하여 자동 스타일링 렌더링
            registry = get_component_registry()
            comp_data = registry.get(comp.component_type, registry.get("Default", {}))
            style = comp_data.get("ppt_style", registry.get("Default", {}).get("ppt_style", {}))
            
            # JSON의 문자열("ROUNDED_RECTANGLE" 등)을 파이썬 MSO_SHAPE 객체로 변환 (오타 시 RECTANGLE 롤백)
            shape_name = style.get("shape", "RECTANGLE")
            if not isinstance(shape_name, str): shape_name = "RECTANGLE"
            shape_enum = getattr(MSO_SHAPE, shape_name, MSO_SHAPE.RECTANGLE)
            
            shape = slide.shapes.add_shape(shape_enum, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = _safe_rgb(style.get("bg"), (245, 245, 255))
            shape.line.color.rgb = _safe_rgb(style.get("line"), (150, 150, 200))
            
            tf = shape.text_frame
            safe_text = comp.text if comp.text else ""
            tf.text = f"[{comp.component_type}]\n{safe_text}".strip()
            tf.word_wrap = True
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            tf.paragraphs[0].font.color.rgb = _safe_rgb(style.get("text"), (100, 100, 150))
            
            try: font_size = int(style.get("size", 12))
            except Exception: font_size = 12
            tf.paragraphs[0].font.size = Pt(font_size)
            
            bold_val = style.get("bold", False)
            if isinstance(bold_val, str):
                tf.paragraphs[0].font.bold = bold_val.strip().lower() in ("true", "1", "yes", "t")
            else:
                tf.paragraphs[0].font.bold = bool(bold_val)

    # --- [우측 화면 설명(Description) 영역 추가 시작] ---
    desc_left = int(actual_slide_width * 0.72)
    desc_top = offset_top
    desc_width = int(actual_slide_width * 0.25)
    desc_height = slide_height
    
    desc_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, desc_left, desc_top, desc_width, desc_height)
    desc_shape.fill.solid()
    desc_shape.fill.fore_color.rgb = RGBColor(250, 250, 250) # 옅은 배경색
    desc_shape.line.color.rgb = RGBColor(200, 200, 200)
    
    tf_desc = desc_shape.text_frame
    tf_desc.word_wrap = True
    tf_desc.margin_left = Inches(0.15)
    tf_desc.margin_right = Inches(0.15)
    tf_desc.margin_top = Inches(0.15)
    tf_desc.margin_bottom = Inches(0.15)
    
    p_title = tf_desc.paragraphs[0]
    p_title.text = "📌 UI 화면 설명 및 퍼블리싱 정책"
    p_title.font.bold = True
    p_title.font.size = Pt(14)
    p_title.font.color.rgb = RGBColor(30, 30, 30)
    
    p_content = tf_desc.add_paragraph()
    p_content.text = f"\n{analysis_result.screen_description}\n\n[컴포넌트별 세부 기능]"
    p_content.font.size = Pt(12)
    p_content.font.color.rgb = RGBColor(80, 80, 80)
    
    for comp in analysis_result.components:
        if comp.event_description:
            p_event = tf_desc.add_paragraph()
            id_prefix = f"[{comp.component_id}] " if hasattr(comp, 'component_id') and comp.component_id else ""
            p_event.text = f"• {id_prefix}{comp.text or comp.component_type}: {comp.event_description}"
            p_event.font.size = Pt(11)
            p_event.font.color.rgb = RGBColor(100, 100, 100)
    # --- [우측 화면 설명(Description) 영역 추가 끝] ---
    
    # --- [하단 날짜 및 페이지 번호 자동 삽입 시작] ---
    current_date = datetime.datetime.now().strftime("%Y-%m-%d")
    slide_number = len(prs.slides) # 현재 생성된 슬라이드의 순번
    
    footer_top = actual_slide_height - Inches(0.4) # 슬라이드 맨 아래에서 약간 위쪽
    
    # 좌측 하단: 오늘 날짜
    date_box = slide.shapes.add_textbox(Inches(0.5), footer_top, Inches(2), Inches(0.3))
    tf_date = date_box.text_frame
    tf_date.text = current_date
    tf_date.paragraphs[0].font.size = Pt(10)
    tf_date.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
    
    # 우측 하단: 페이지 번호
    page_box = slide.shapes.add_textbox(actual_slide_width - Inches(2.5), footer_top, Inches(2), Inches(0.3))
    tf_page = page_box.text_frame
    tf_page.text = f"- {slide_number} -"
    tf_page.paragraphs[0].alignment = PP_ALIGN.RIGHT
    tf_page.paragraphs[0].font.size = Pt(10)
    tf_page.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
    # --- [하단 날짜 및 페이지 번호 자동 삽입 끝] ---

    # 지정된 경로에 파일 저장
    prs.save(output_file)

def sync_design_tokens_from_ppt(template_path):
    """PPT 템플릿의 도형을 스캔하여 components_registry.json의 ppt_style을 자동 추출/업데이트합니다."""
    if not template_path or not os.path.exists(template_path):
        return 0, "PPT 파일을 찾을 수 없습니다."
        
    try:
        prs = Presentation(template_path)
        registry = get_component_registry()
        updated_count = 0
        
        for slide in prs.slides:
            for shape in slide.shapes:
                comp_name = None
                
                # 1. PPT 도형의 고유 '이름(Name)' 속성 확인 (선택 창에서 변경한 커스텀 이름)
                # 영문자, 숫자, 하이픈(-), 언더바(_) 조합인지 정규식으로 유연하게 확인
                if shape.name and re.match(r'^[A-Za-z0-9_-]{3,}$', shape.name):
                    comp_name = shape.name
                    
                # 2. 도형 내부 텍스트 확인 ([HdsButton] 또는 HdsButton)
                if not comp_name and shape.has_text_frame:
                    text = shape.text.strip()
                    if text.startswith("[") and "]" in text:
                        comp_name = text[1:text.find("]")]
                    elif text and re.match(r'^[A-Za-z0-9_-]{3,}$', text):
                        comp_name = text
                        
                if not comp_name:
                    continue
                
                # 1. 형태(Shape) 스캔
                shape_str = "RECTANGLE"
                if hasattr(shape, "auto_shape_type") and shape.auto_shape_type:
                    # MSO_SHAPE 열거형 값(정수)을 다시 문자열 이름("ROUNDED_RECTANGLE" 등)으로 역변환
                    for key in dir(MSO_SHAPE):
                        if not key.startswith("_") and getattr(MSO_SHAPE, key) == shape.auto_shape_type:
                            shape_str = key
                            break
                    
                # 2. 배경색(Fill) 스캔
                bg_color = [245, 245, 255]
                try:
                    if hasattr(shape.fill, "fore_color") and hasattr(shape.fill.fore_color, "rgb") and shape.fill.fore_color.rgb:
                        rgb = shape.fill.fore_color.rgb
                        bg_color = [rgb[0], rgb[1], rgb[2]]
                except Exception: pass
                
                # 3. 테두리색(Line) 스캔
                line_color = [200, 200, 200]
                try:
                    if hasattr(shape.line, "color") and hasattr(shape.line.color, "rgb") and shape.line.color.rgb:
                        rgb = shape.line.color.rgb
                        line_color = [rgb[0], rgb[1], rgb[2]]
                except Exception: pass
                
                # 4. 텍스트 색상 및 폰트 속성 스캔
                text_color = [50, 50, 50]
                font_size = 12
                font_bold = False
                try:
                    if shape.has_text_frame and shape.text_frame.paragraphs:
                        font = shape.text_frame.paragraphs[0].font
                        if hasattr(font, "color") and hasattr(font.color, "rgb") and font.color.rgb:
                            rgb = font.color.rgb
                            text_color = [rgb[0], rgb[1], rgb[2]]
                        if font.size:
                            font_size = int(font.size.pt)
                        if font.bold is not None:
                            font_bold = font.bold
                except Exception: pass
                
                # 5. 레지스트리 병합(Merge) 저장
                if comp_name in registry and registry[comp_name].get("locked") is True:
                    continue # 잠금(Lock) 설정된 컴포넌트는 업데이트 스킵
                    
                if comp_name not in registry:
                    registry[comp_name] = {"guide": "", "ppt_style": {}}
                    
                registry[comp_name]["ppt_style"] = {
                    "shape": shape_str, "bg": bg_color, "line": line_color, 
                    "text": text_color, "size": font_size, "bold": font_bold
                }
                updated_count += 1
                    
        # 추출된 데이터가 있으면 JSON 파일 덮어쓰기
        if updated_count > 0:
            with open("components_registry.json", "w", encoding="utf-8") as f:
                json.dump(registry, f, ensure_ascii=False, indent=4)
        return updated_count, "성공"
    except Exception as e:
        logger.error(f"디자인 토큰 스캔 중 오류: {e}", exc_info=True)
        return 0, str(e)

def sync_design_tokens_from_css(css_path="hds.css"):
    """CSS 파일의 속성을 파싱하여 components_registry.json의 ppt_style을 자동 업데이트합니다."""
    if not os.path.exists(css_path):
        return 0, "hds.css 파일을 찾을 수 없습니다."
        
    try:
        with open(css_path, "r", encoding="utf-8") as f:
            css_content = f.read()
            
        registry = get_component_registry()
        updated_count = 0
        
        # CSS 블록 추출 (예: .ant-btn-primary { background-color: #E60012; })
        blocks = re.findall(r'\.([a-zA-Z0-9_-]+)\s*\{([^}]*)\}', css_content)
        
        for class_name, block_content in blocks:
            bg_match = re.search(r'background(?:-color)?\s*:\s*([^;!]+)', block_content)
            text_match = re.search(r'(?<!-|\w)color\s*:\s*([^;!]+)', block_content)
            border_match = re.search(r'border(?:-color)?\s*:\s*([^;!]+)', block_content)
            
            bg_color = _parse_color(bg_match.group(1)) if bg_match else None
            text_color = _parse_color(text_match.group(1)) if text_match else None
            line_color = None
            if border_match:
                hex_in_border = re.search(r'#[0-9a-fA-F]+', border_match.group(1))
                if hex_in_border: line_color = _hex_to_rgb(hex_in_border.group(0))
                else: line_color = _parse_color(border_match.group(1))
                
            if bg_color or text_color or line_color:
                if class_name in registry and registry[class_name].get("locked") is True:
                    continue # 잠금(Lock) 설정된 컴포넌트는 업데이트 스킵
                    
                if class_name not in registry:
                    registry[class_name] = {"guide": f"CSS 클래스 .{class_name} 기반 컴포넌트", "ppt_style": {}}
                
                ppt_style = registry[class_name].get("ppt_style", {})
                if "shape" not in ppt_style:
                    ppt_style["shape"] = "ROUNDED_RECTANGLE" if "btn" in class_name else "RECTANGLE"
                    
                if bg_color: ppt_style["bg"] = bg_color
                if text_color: ppt_style["text"] = text_color
                if line_color: ppt_style["line"] = line_color
                
                registry[class_name]["ppt_style"] = ppt_style
                updated_count += 1
                
        if updated_count > 0:
            with open("components_registry.json", "w", encoding="utf-8") as f:
                json.dump(registry, f, ensure_ascii=False, indent=4)
                
        return updated_count, "성공"
    except Exception as e:
        logger.error(f"CSS 스캔 중 오류: {e}", exc_info=True)
        return 0, str(e)